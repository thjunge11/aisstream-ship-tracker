"""
Generates AIS Ship Tracker – Capstone Presentation.pptx
from the content of presentation.html (converted manually to python-pptx calls).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Colour palette (mirrors the HTML CSS) ───────────────────────────────────
DARK       = RGBColor(0x0D, 0x11, 0x17)   # slide background
ACCENT     = RGBColor(0x00, 0xC8, 0xFF)   # cyan
ACCENT2    = RGBColor(0x00, 0xE6, 0x76)   # green
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREY       = RGBColor(0xAA, 0xAA, 0xAA)
ORANGE     = RGBColor(0xFF, 0xA5, 0x00)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)
CODE_BG    = RGBColor(0x16, 0x1B, 0x22)
CODE_FG    = RGBColor(0xC9, 0xD1, 0xD9)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    slide = prs.slides.add_slide(blank_layout)
    # Dark background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK
    return slide


def add_textbox(slide, left, top, width, height,
                text, font_size=18, bold=False, italic=False,
                color=WHITE, align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_para(tf, text, font_size=14, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, bullet=False):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.alignment = align
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = ("• " if bullet else "") + text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def section_label(slide, text, top=Inches(0.18)):
    add_textbox(slide, Inches(0.4), top, Inches(12), Inches(0.3),
                text, font_size=9, color=RGBColor(0x55, 0x55, 0x55))


def slide_title(slide, text, top=Inches(0.4)):
    txb = slide.shapes.add_textbox(Inches(0.4), top, Inches(12.5), Inches(0.55))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    return txb


def add_card(slide, left, top, width, height,
             title=None, bullets=None, font_size=11):
    """Draw a bordered card with optional title and bullet list."""
    # border rectangle
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x20, 0x28)
    shape.line.color.rgb = RGBColor(0x00, 0x64, 0x88)
    shape.line.width = Pt(0.75)

    inner_left  = left  + Inches(0.1)
    inner_top   = top   + Inches(0.08)
    inner_width = width - Inches(0.2)

    if title:
        add_textbox(slide, inner_left, inner_top, inner_width, Inches(0.3),
                    title, font_size=font_size + 1, bold=True, color=ACCENT2)
        inner_top += Inches(0.28)

    if bullets:
        txb = slide.shapes.add_textbox(inner_left, inner_top,
                                        inner_width,
                                        height - (inner_top - top) - Inches(0.05))
        txb.word_wrap = True
        tf = txb.text_frame
        tf.word_wrap = True
        first = True
        for item in bullets:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(1)
            run = p.add_run()
            run.text = "• " + item
            run.font.size = Pt(font_size)
            run.font.color.rgb = GREY

    return shape


def add_pill(slide, left, top, text, color=ACCENT):
    """Rounded rectangle pill badge."""
    w = max(Inches(1.0), Inches(len(text) * 0.09 + 0.3))
    h = Inches(0.27)
    shape = slide.shapes.add_shape(5, left, top, w, h)  # 5 = rounded rect
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x20, 0x28)
    shape.line.color.rgb = color
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.margin_left = util.Pt(4)
    tf.margin_right = util.Pt(4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = color
    return shape, w


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ════════════════════════════════════════════════════════════════════════════

def slide_01(prs):
    slide = blank_slide(prs)

    # Anchor emoji
    add_textbox(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(1.0),
                "⚓", font_size=54, align=PP_ALIGN.CENTER)

    # Main title
    add_textbox(slide, Inches(0.5), Inches(2.3), Inches(12.33), Inches(0.9),
                "AIS Ship Tracker", font_size=44, bold=True,
                color=ACCENT, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(3.15), Inches(12.33), Inches(0.4),
                "A real-time maritime data engineering platform",
                font_size=16, color=GREY, align=PP_ALIGN.CENTER)

    # Tagline
    add_textbox(slide, Inches(0.5), Inches(3.6), Inches(12.33), Inches(0.35),
                "DATA ENGINEERING CAPSTONE PROJECT  ·  MAY 2026",
                font_size=11, color=ACCENT2, align=PP_ALIGN.CENTER)

    # Badges row
    badges = ["Kafka", "Python", "PostgreSQL", "Docker",
              "Flask", "MapLibre GL", "AWS EC2", "GitHub Actions"]
    total_w = sum(max(Inches(1.0), Inches(len(b) * 0.09 + 0.3)) for b in badges)
    gap = Inches(0.12)
    start_x = (SLIDE_W - total_w - gap * (len(badges) - 1)) / 2
    x = start_x
    y = Inches(4.15)
    for badge in badges:
        shape, w = add_pill(slide, x, y, badge)
        x += w + gap


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – BACKGROUND
# ════════════════════════════════════════════════════════════════════════════

def slide_02(prs):
    slide = blank_slide(prs)
    section_label(slide, "Background")
    slide_title(slide, "What is AIS & Why does it matter?")

    # Left column
    add_textbox(slide, Inches(0.4), Inches(1.05), Inches(6.1), Inches(0.28),
                "The Problem", font_size=12, bold=True, color=ACCENT2)

    left_tb = slide.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(6.1), Inches(2.4))
    left_tb.word_wrap = True
    tf = left_tb.text_frame
    tf.word_wrap = True
    bullets_problem = [
        "Over 90 000 vessels are at sea at any moment",
        "The Automatic Identification System (AIS) broadcasts GPS position, speed, heading and identity every few seconds from every vessel",
        "The raw stream is a high-velocity, schema-varied firehose — unsuitable for direct consumption",
        "No single open tool turns the raw WebSocket feed into a queryable, live-updating map",
    ]
    first = True
    for b in bullets_problem:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(11)
        run.font.color.rgb = GREY

    add_textbox(slide, Inches(0.4), Inches(3.85), Inches(6.1), Inches(0.28),
                "The Goal", font_size=12, bold=True, color=ACCENT2)

    goal_tb = slide.shapes.add_textbox(Inches(0.4), Inches(4.15), Inches(6.1), Inches(2.2))
    goal_tb.word_wrap = True
    tf2 = goal_tb.text_frame
    tf2.word_wrap = True
    bullets_goal = [
        "Ingest the global AIS WebSocket stream continuously",
        "Route, validate and persist positions in real-time",
        "Surface live positions and historical tracks on an interactive map",
        "Deploy the whole stack to the cloud with zero-touch CI/CD",
    ]
    first = True
    for b in bullets_goal:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        run = p.add_run()
        run.text = "• " + b
        run.font.size = Pt(11)
        run.font.color.rgb = GREY

    # Right column – two cards
    add_card(slide, Inches(6.8), Inches(1.0), Inches(6.1), Inches(2.1),
             title="Data Source – AISStream.io",
             bullets=[
                 "WebSocket API at wss://stream.aisstream.io/v0/stream",
                 "Subscription filter: global bounding box [-90,-180] → [90,180]",
                 "Each message carries a MessageType field — 20+ distinct AIS message types are fanned out automatically to matching Kafka topics.",
             ], font_size=11)

    add_card(slide, Inches(6.8), Inches(3.25), Inches(6.1), Inches(2.0),
             title="Key AIS Message Types",
             bullets=[
                 "PositionReport – lat/lon, speed, heading (every ~5 s)",
                 "ShipStaticData – name, MMSI, type, dimensions",
                 "StandardClassBPositionReport, StaticDataReport, …",
             ], font_size=11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════

def slide_03(prs):
    slide = blank_slide(prs)
    section_label(slide, "System Design")
    slide_title(slide, "Architecture Overview")

    # Pipeline row
    pipeline = [
        ("AISStream.io\nWebSocket", ACCENT),
        ("ais_producer\nPython", ACCENT),
        ("Apache Kafka\n20+ topics", ACCENT),
        ("live_data_processor\nPython", ACCENT),
        ("ships_live_data\nKafka topic", ACCENT2),
        ("PostgreSQL\n3 tables", ORANGE),
        ("Flask + MapLibre\nWeb App", ACCENT2),
    ]
    x = Inches(0.3)
    y = Inches(1.05)
    for i, (label, col) in enumerate(pipeline):
        w = Inches(1.7)
        h = Inches(0.55)
        shape = slide.shapes.add_shape(5, x, y, w, h)
        shape.adjustments[0] = 0.15
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1A, 0x20, 0x28)
        shape.line.color.rgb = col
        shape.line.width = Pt(0.75)
        tf = shape.text_frame
        tf.margin_left = Pt(4)
        tf.margin_right = Pt(4)
        for j, line in enumerate(label.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(8)
            run.font.color.rgb = col
        x += w
        if i < len(pipeline) - 1:
            add_textbox(slide, x, y + Inches(0.15), Inches(0.2), Inches(0.28),
                        "→", font_size=14, color=GREY)
            x += Inches(0.2)

    # Four cards 2x2
    cards = [
        ("Ingestion Layer", [
            "ais_producer – async WebSocket consumer; fans out each message to the Kafka topic matching its MessageType",
            "Emits throughput stats every 10 000 messages to statistics_ais_producer",
        ]),
        ("Processing Layer", [
            "ships_live_data_processor – validates & transforms PositionReport → ships_live_data topic",
            "ships_static_data_processor – transforms ShipStaticData → ships_static_data topic",
        ]),
        ("Persistence Layer", [
            "ships_live_data_consumer – upserts latest position into ships_live_data table",
            "position_history_consumer – writes timestamped positions for tracked ships into position_history",
            "ships_static_data_consumer – upserts static vessel info",
        ]),
        ("Serving Layer", [
            "Flask REST API – /api/ships/live, /api/ships/<id>/history, /api/ships/tracked",
            "MapLibre GL JS – renders live ship arrows on a dark tile map",
        ]),
    ]
    positions = [
        (Inches(0.3),  Inches(1.75)),
        (Inches(6.8),  Inches(1.75)),
        (Inches(0.3),  Inches(4.2)),
        (Inches(6.8),  Inches(4.2)),
    ]
    for (left, top), (title, bullets) in zip(positions, cards):
        add_card(slide, left, top, Inches(6.1), Inches(2.2),
                 title=title, bullets=bullets, font_size=10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – KAFKA TOPICS
# ════════════════════════════════════════════════════════════════════════════

def slide_04(prs):
    slide = blank_slide(prs)
    section_label(slide, "Message Broker")
    slide_title(slide, "Apache Kafka – Topic Landscape")

    # Left column intro
    add_textbox(slide, Inches(0.4), Inches(1.05), Inches(6.1), Inches(0.8),
                "Every AIS message type becomes its own Kafka topic. The producer routes by message[\"MessageType\"] with zero configuration — new types are automatically handled.",
                font_size=11, color=GREY)

    add_card(slide, Inches(0.4), Inches(1.95), Inches(6.1), Inches(1.55),
             title="Kafka Setup",
             bullets=[
                 "Single broker · Zookeeper · Confluent Platform 7.5.3",
                 "All topics: 1 partition, RF 1 (single-node)",
                 "Configurable retention via KAFKA_LOG_RETENTION_MS",
                 "Kafka UI (Provectus) for observability on port 8080",
             ], font_size=10)

    # High-volume topics table
    add_textbox(slide, Inches(0.4), Inches(3.62), Inches(6.1), Inches(0.28),
                "High-volume topics (observed)", font_size=11, bold=True, color=ACCENT2)

    topics = [
        ("PositionReport",                   "239 073 msgs", "153 MB"),
        ("StandardClassBPositionReport",      "76 264",       "58 MB"),
        ("ShipStaticData",                    "42 701",       "28 MB"),
        ("DataLinkManagementMessage",         "30 929",       "23 MB"),
        ("StaticDataReport",                  "32 560",       "21 MB"),
    ]
    row_h = Inches(0.3)
    for i, (name, count, size) in enumerate(topics):
        y = Inches(3.95) + i * row_h
        add_textbox(slide, Inches(0.4),  y, Inches(3.1), row_h, name,  font_size=10, color=ACCENT)
        add_textbox(slide, Inches(3.55), y, Inches(1.5), row_h, count, font_size=10, color=ACCENT2)
        add_textbox(slide, Inches(5.1),  y, Inches(1.3), row_h, size,  font_size=10, color=GREY)

    # Right column – stats JSON (as code-style box)
    add_textbox(slide, Inches(6.8), Inches(1.05), Inches(6.1), Inches(0.28),
                "Producer throughput statistics", font_size=11, bold=True, color=ACCENT2)

    code_text = (
        '// emitted every 10 000 messages\n'
        '// topic: statistics_ais_producer\n'
        '{\n'
        '  "timestamp":             "2026-04-25T10:20:50+00:00",\n'
        '  "interval_seconds":       37.9,\n'
        '  "interval_message_count": 10000,\n'
        '  "interval_message_rate":  263.8,   // msgs/s\n'
        '  "message_count":          380000\n'
        '}'
    )
    code_box = slide.shapes.add_shape(1, Inches(6.8), Inches(1.37), Inches(6.1), Inches(2.15))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = CODE_BG
    code_box.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    add_textbox(slide, Inches(6.9), Inches(1.42), Inches(5.9), Inches(2.05),
                code_text, font_size=9, color=CODE_FG)

    add_card(slide, Inches(6.8), Inches(3.62), Inches(6.1), Inches(1.1),
             title="Throughput",
             bullets=[
                 "Sustained ~250 msg/s from the global AIS stream",
                 "720 000+ messages processed in a single session",
                 "At peak: one position update every 4 ms",
             ], font_size=10)

    add_card(slide, Inches(6.8), Inches(4.82), Inches(6.1), Inches(1.55),
             title="Message routing (producer core)",
             bullets=[
                 'async for message_json in websocket:',
                 '    message = json.loads(message_json)',
                 '    producer.send(message["MessageType"], message)',
             ], font_size=10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – STREAM PROCESSING
# ════════════════════════════════════════════════════════════════════════════

def slide_05(prs):
    slide = blank_slide(prs)
    section_label(slide, "Python Microservices")
    slide_title(slide, "Stream Processing Pipeline")

    services = [
        ("ships_live_data_processor",
         "In: PositionReport  →  Out: ships_live_data",
         [
             "Validates AIS sentinel values (lat 91°, lon 181°, COG 360°, SOG 102.3+)",
             "Maps numeric nav-status code → human-readable string (ITU-R M.1371-5)",
             "Emits clean, typed JSON with ISO-8601 UTC timestamp",
             "Skipped messages logged with reason",
         ]),
        ("ships_live_data_consumer",
         "In: ships_live_data  →  Out: PostgreSQL",
         [
             "Batched upsert into ships_live_data table (configurable batch size)",
             "ON CONFLICT (ship_id) DO UPDATE keeps only latest position",
             "Table stays small – one row per ship regardless of volume",
         ]),
        ("position_history_consumer",
         "In: ships_live_data  →  Out: PostgreSQL",
         [
             "Only writes rows for ships listed in tracking_config",
             "Tracking config refreshed every N messages (no restart required)",
             "35+ batches written; 139 000+ history points accumulated for 6 ships",
         ]),
        ("ships_static_data_processor",
         "In: ShipStaticData  →  Out: ships_static_data",
         [
             "Extracts name, call sign, IMO number, vessel type, dimensions, ETA",
             "Joins ship identity to live data via shared MMSI (ship_id)",
         ]),
        ("ships_static_data_consumer",
         "In: ships_static_data  →  Out: PostgreSQL",
         [
             "Batched upsert into ships_static_data table",
             "Configurable batch size via env var",
         ]),
        ("tracking_sync",
         "Trigger: GitHub Actions",
         [
             "Reads tracked_ship_ids.json from the repo",
             "Reconciles against tracking_config in Postgres",
             "Activates / deactivates tracking without app restart",
             "Ships not yet seen in the stream are skipped and retried on next run",
         ]),
    ]

    # 3 columns × 2 rows
    col_w = Inches(4.2)
    col_gap = Inches(0.2)
    row_h = Inches(2.7)
    for i, (name, io_line, bullets) in enumerate(services):
        col = i % 3
        row = i // 3
        left = Inches(0.3) + col * (col_w + col_gap)
        top  = Inches(1.05) + row * (row_h + Inches(0.15))
        add_card(slide, left, top, col_w, row_h,
                 title=name,
                 bullets=[io_line, ""] + bullets,
                 font_size=10)
        # Override the IO line colour — just write it on top
        add_textbox(slide, left + Inches(0.1), top + Inches(0.35),
                    col_w - Inches(0.2), Inches(0.25),
                    io_line, font_size=9, color=GREY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – DATA MODEL
# ════════════════════════════════════════════════════════════════════════════

def slide_06(prs):
    slide = blank_slide(prs)
    section_label(slide, "Storage")
    slide_title(slide, "PostgreSQL Data Model")

    def schema_block(slide, left, top, width, title, rows, font_size=10):
        add_textbox(slide, left, top, width, Inches(0.28),
                    title, font_size=font_size + 1, bold=True, color=ACCENT2)
        top += Inches(0.3)
        hdr_h = Inches(0.26)
        shape = slide.shapes.add_shape(1, left, top, width, hdr_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x00, 0x40, 0x55)
        shape.line.color.rgb = RGBColor(0x00, 0x64, 0x88)
        for j, hdr in enumerate(["Column", "Type", ""]):
            x = left + j * (width / 3)
            add_textbox(slide, x, top, width / 3, hdr_h,
                        hdr, font_size=font_size - 1, bold=True, color=ACCENT)
        top += hdr_h
        for col_name, col_type, note in rows:
            row_shape = slide.shapes.add_shape(1, left, top, width, hdr_h)
            row_shape.fill.solid()
            row_shape.fill.fore_color.rgb = RGBColor(0x1A, 0x20, 0x28)
            row_shape.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
            add_textbox(slide, left,                top, width / 3, hdr_h, col_name,  font_size=font_size - 1, color=GREY)
            add_textbox(slide, left + width / 3,    top, width / 3, hdr_h, col_type,  font_size=font_size - 1, color=RGBColor(0xA5, 0xD6, 0xFF))
            note_color = GOLD if "PK" in note else (RGBColor(0x87, 0xCE, 0xEB) if "FK" in note else GREY)
            add_textbox(slide, left + 2 * width / 3, top, width / 3, hdr_h, note,    font_size=font_size - 1, color=note_color)
            top += hdr_h
        return top

    # Left column
    schema_block(slide, Inches(0.4), Inches(1.0), Inches(6.1),
                 "ships_live_data — latest position per ship",
                 [
                     ("ship_id",              "INTEGER",          "PK"),
                     ("ship_name",            "VARCHAR(50)",      ""),
                     ("course_over_ground",   "FLOAT",            ""),
                     ("speed_over_ground",    "FLOAT",            ""),
                     ("navigational_status",  "VARCHAR(100)",     ""),
                     ("latitude / longitude", "DOUBLE PRECISION", ""),
                     ("updated_at",           "TIMESTAMPTZ",      ""),
                 ])

    schema_block(slide, Inches(0.4), Inches(4.15), Inches(6.1),
                 "tracking_config — opt-in history",
                 [
                     ("ship_id",      "INTEGER",   "PK FK"),
                     ("enabled_from", "TIMESTAMPTZ", "PK"),
                     ("enabled_to",   "TIMESTAMPTZ", "NULL = active"),
                 ])

    # Right column
    schema_block(slide, Inches(6.8), Inches(1.0), Inches(6.1),
                 "position_history — time-series track",
                 [
                     ("id",                   "BIGSERIAL",        "PK"),
                     ("ship_id",              "INTEGER",          "FK"),
                     ("course_over_ground",   "FLOAT",            ""),
                     ("speed_over_ground",    "FLOAT",            ""),
                     ("navigational_status",  "VARCHAR(100)",     ""),
                     ("latitude / longitude", "DOUBLE PRECISION", ""),
                     ("recorded_at",          "TIMESTAMPTZ",      ""),
                 ])

    schema_block(slide, Inches(6.8), Inches(3.95), Inches(6.1),
                 "ships_static_data",
                 [
                     ("ship_id",                   "INTEGER PK",  ""),
                     ("ship_name, call_sign, IMO",  "VARCHAR",     ""),
                     ("ship_type, destination, ETA","VARCHAR",     ""),
                     ("length_m, width_m, draught", "INT / FLOAT", ""),
                     ("updated_at",                 "TIMESTAMPTZ", ""),
                 ])

    add_card(slide, Inches(6.8), Inches(5.8), Inches(6.1), Inches(1.55),
             title="Design decisions",
             bullets=[
                 "ships_live_data is an upsert-only table — always O(1) lookup for the map",
                 "tracking_config uses composite PK (ship_id, enabled_from) to support multiple tracking periods",
                 "FK cascade ensures history is always tied to a known vessel",
             ], font_size=10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – WEB APP
# ════════════════════════════════════════════════════════════════════════════

def slide_07(prs):
    slide = blank_slide(prs)
    section_label(slide, "Frontend")
    slide_title(slide, "Web Application – Live Ship Map")

    add_card(slide, Inches(0.4), Inches(1.0), Inches(6.1), Inches(1.55),
             title="Stack",
             bullets=[
                 "Flask (Python) – lightweight REST API server",
                 "MapLibre GL JS – GPU-accelerated WebGL map rendering",
                 "GeoJSON – all spatial data transferred as standard FeatureCollections",
                 "Vanilla JS – no frontend framework overhead",
             ], font_size=11)

    add_card(slide, Inches(0.4), Inches(2.65), Inches(6.1), Inches(1.75),
             title="API Endpoints",
             bullets=[
                 "GET /api/ships/live  →  GeoJSON of all ships with lat/lon",
                 "GET /api/ships/<id>/history  →  GeoJSON LineString + Points (up to 200 positions)",
                 "GET /api/ships/tracked  →  JSON list of actively-tracked ship IDs",
             ], font_size=11)

    add_card(slide, Inches(0.4), Inches(4.5), Inches(6.1), Inches(1.9),
             title="Live view stats (observed)",
             bullets=[
                 "26 100 ships visible simultaneously on the map",
                 "Ship arrows rendered with MapLibre symbol layer, rotated by course_over_ground",
                 "Tracked ships highlighted in a distinct colour with clickable history track",
                 "Sidebar search by ship name or MMSI",
             ], font_size=11)

    # Right column — code blocks
    add_textbox(slide, Inches(6.8), Inches(1.0), Inches(6.1), Inches(0.28),
                "No-cache API pattern", font_size=11, bold=True, color=ACCENT2)

    code1 = (
        "@app.after_request\n"
        "def no_cache(response):\n"
        '    if request.path.startswith("/api/"):\n'
        '        response.headers["Cache-Control"] = \\\n'
        '            "no-store, no-cache"\n'
        "    return response"
    )
    box1 = slide.shapes.add_shape(1, Inches(6.8), Inches(1.3), Inches(6.1), Inches(1.5))
    box1.fill.solid(); box1.fill.fore_color.rgb = CODE_BG
    box1.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    add_textbox(slide, Inches(6.9), Inches(1.35), Inches(5.9), Inches(1.4),
                code1, font_size=9.5, color=CODE_FG)

    add_textbox(slide, Inches(6.8), Inches(2.9), Inches(6.1), Inches(0.28),
                "GeoJSON live response shape", font_size=11, bold=True, color=ACCENT2)

    code2 = (
        '{\n'
        '  "type": "FeatureCollection",\n'
        '  "features": [{\n'
        '    "type": "Feature",\n'
        '    "geometry": {\n'
        '      "type": "Point",\n'
        '      "coordinates": [10.12, 54.33]\n'
        '    },\n'
        '    "properties": {\n'
        '      "ship_id": 211341930,\n'
        '      "ship_name": "GAARDEN",\n'
        '      "speed_over_ground": 1.0,\n'
        '      "course_over_ground": 17.1,\n'
        '      "is_tracked": true\n'
        '    }\n'
        '  }]\n'
        '}'
    )
    box2 = slide.shapes.add_shape(1, Inches(6.8), Inches(3.2), Inches(6.1), Inches(3.15))
    box2.fill.solid(); box2.fill.fore_color.rgb = CODE_BG
    box2.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    add_textbox(slide, Inches(6.9), Inches(3.25), Inches(5.9), Inches(3.05),
                code2, font_size=9.5, color=CODE_FG)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – DEVOPS
# ════════════════════════════════════════════════════════════════════════════

def slide_08(prs):
    slide = blank_slide(prs)
    section_label(slide, "DevOps & Infrastructure")
    slide_title(slide, "Docker, CI/CD & AWS Deployment")

    # Left column — container stack
    add_textbox(slide, Inches(0.4), Inches(1.05), Inches(6.1), Inches(0.28),
                "Container Stack", font_size=12, bold=True, color=ACCENT2)
    add_textbox(slide, Inches(0.4), Inches(1.35), Inches(6.1), Inches(0.5),
                "Every service runs in its own Docker container, coordinated by a single docker-compose.yml:",
                font_size=10, color=GREY)

    pill_groups = [
        (["zookeeper", "kafka", "kafka-ui", "postgres"], ACCENT),
        (["ais-producer", "elt-processor", "live-data-consumer"], ACCENT2),
        (["position-history-consumer", "static-data-processor", "static-data-consumer"], ACCENT2),
        (["webapp", "tracking-sync"], ORANGE),
    ]
    y = Inches(1.9)
    for group, col in pill_groups:
        x = Inches(0.4)
        for label in group:
            shape, w = add_pill(slide, x, y, label, color=col)
            x += w + Inches(0.1)
        y += Inches(0.36)

    # GitHub Actions workflows
    add_textbox(slide, Inches(0.4), Inches(3.45), Inches(6.1), Inches(0.28),
                "GitHub Actions Workflows", font_size=12, bold=True, color=ACCENT2)

    workflows = [
        ("provision-ec2.yml", "Creates EC2 instance, security group, bootstraps Docker — manual trigger"),
        ("deploy.yml",        "SSH deploy on PR merge to main — git pull + docker compose up"),
        ("tracking-sync.yml", "Runs on change to tracked_ship_ids.json — syncs DB tracking config"),
    ]
    wy = Inches(3.78)
    for wf_name, wf_desc in workflows:
        shape = slide.shapes.add_shape(1, Inches(0.4), wy, Inches(2.1), Inches(0.42))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x00, 0x30, 0x44)
        shape.line.color.rgb = ACCENT
        add_textbox(slide, Inches(0.45), wy + Inches(0.05), Inches(2.0), Inches(0.38),
                    wf_name, font_size=9, color=ACCENT)
        add_textbox(slide, Inches(2.6), wy + Inches(0.05), Inches(3.85), Inches(0.38),
                    wf_desc, font_size=9, color=GREY)
        wy += Inches(0.5)

    # Right column
    add_card(slide, Inches(6.8), Inches(1.0), Inches(6.1), Inches(2.05),
             title="AWS Infrastructure",
             bullets=[
                 "Single t3.large EC2 instance (2 vCPU / 8 GB RAM)",
                 "Ubuntu 22.04 LTS — resolved dynamically via SSM parameter (no stale AMI IDs)",
                 "30 GB gp3 EBS volume",
                 "Security group: ports 22 (SSH), 5000 (webapp), 9093 (Kafka)",
                 "Tagged ManagedBy: github-actions",
             ], font_size=11)

    add_card(slide, Inches(6.8), Inches(3.15), Inches(6.1), Inches(1.55),
             title="Secrets Management",
             bullets=[
                 "All credentials stored as GitHub repository secrets",
                 "ENV_FILE secret is written to .env on the instance via printf — never touches git history",
                 "API key, DB password, AIS key — zero secrets in code",
             ], font_size=11)

    add_textbox(slide, Inches(6.8), Inches(4.8), Inches(6.1), Inches(0.28),
                "Deploy flow", font_size=11, bold=True, color=ACCENT2)

    deploy_code = (
        "# On every PR merge → main\n"
        "git pull origin main\n"
        "printf '%s' \"$ENV_FILE\" > .env.tmp\n"
        "mv .env.tmp .env\n"
        "docker compose up --build -d \\\n"
        "    --remove-orphans\n"
        "docker image prune -f"
    )
    box = slide.shapes.add_shape(1, Inches(6.8), Inches(5.1), Inches(6.1), Inches(1.85))
    box.fill.solid(); box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    add_textbox(slide, Inches(6.9), Inches(5.15), Inches(5.9), Inches(1.75),
                deploy_code, font_size=9.5, color=CODE_FG)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – METRICS
# ════════════════════════════════════════════════════════════════════════════

def slide_09(prs):
    slide = blank_slide(prs)
    section_label(slide, "Results")
    slide_title(slide, "Key Numbers & Performance")

    metrics = [
        ("26 100",    "Ships tracked live simultaneously"),
        ("~264",      "Messages per second sustained"),
        ("720 000+",  "Messages processed in one session"),
        ("20+",       "Distinct Kafka topics auto-routed"),
        ("239 K",     "PositionReport messages (153 MB)"),
        ("139 K+",    "Position history points written"),
        ("11",        "Docker containers in the stack"),
        ("3",         "GitHub Actions workflows"),
    ]

    card_w = Inches(3.0)
    card_h = Inches(1.35)
    gap    = Inches(0.2)
    cols   = 4
    start_x = Inches(0.35)
    start_y = Inches(1.0)

    for i, (value, label) in enumerate(metrics):
        col = i % cols
        row = i // cols
        left = start_x + col * (card_w + gap)
        top  = start_y + row * (card_h + gap)

        shape = slide.shapes.add_shape(1, left, top, card_w, card_h)
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x1A, 0x20, 0x28)
        shape.line.color.rgb = RGBColor(0x00, 0x64, 0x88)
        shape.line.width = Pt(0.75)

        add_textbox(slide, left, top + Inches(0.15), card_w, Inches(0.55),
                    value, font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, left, top + Inches(0.72), card_w, Inches(0.55),
                    label, font_size=10, color=GREY, align=PP_ALIGN.CENTER)

    add_card(slide, Inches(0.35), Inches(4.35), Inches(6.3), Inches(1.6),
             title="Observed throughput log",
             bullets=[
                 "The producer processed 10 000 messages every ~38 seconds on average.",
                 "At this rate the system ingests ~920 000 messages per hour.",
                 "Well within the capacity of a single Kafka broker and single-node Postgres on a t3.large.",
             ], font_size=11)

    add_card(slide, Inches(6.85), Inches(4.35), Inches(6.1), Inches(1.6),
             title="Position history consumer",
             bullets=[
                 "With 6 actively-tracked ships, the position history consumer wrote 35+ batches.",
                 "~139 000 non-tracked messages skipped.",
                 "Full history written for named vessels: COLOR MAGIC, STENA SCANDINAVICA, GAARDEN.",
             ], font_size=11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – ROADMAP
# ════════════════════════════════════════════════════════════════════════════

def slide_10(prs):
    slide = blank_slide(prs)
    section_label(slide, "What's Next")
    slide_title(slide, "Roadmap & Future Work")

    # Left column – Completed & Next
    add_textbox(slide, Inches(0.4), Inches(1.0), Inches(6.1), Inches(0.28),
                "✓ Completed", font_size=12, bold=True, color=ACCENT2)

    done = [
        "Global AIS WebSocket ingestion via aisstream.io",
        "Kafka fan-out by message type (20+ topics)",
        "Stream validation & transformation pipeline",
        "PostgreSQL persistence with live + history tables",
        "Opt-in position history tracking via JSON config",
        "MapLibre GL live map with 26 000+ ships",
        "Full Docker Compose stack",
        "GitHub Actions: provision EC2, deploy on PR merge, tracking sync",
    ]
    y = Inches(1.32)
    for item in done:
        shape = slide.shapes.add_shape(9, Inches(0.4), y + Inches(0.04), Inches(0.14), Inches(0.14))
        shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT2
        shape.line.fill.background()
        add_textbox(slide, Inches(0.6), y, Inches(5.85), Inches(0.3),
                    item, font_size=10, color=GREY)
        y += Inches(0.31)

    add_textbox(slide, Inches(0.4), y + Inches(0.05), Inches(6.1), Inches(0.28),
                "→ Next", font_size=12, bold=True, color=ACCENT)
    y += Inches(0.38)

    nxt = [
        "Ship static data enrichment in the UI (name, type, dimensions)",
        "History track visualisation on the map (LineString layer)",
        "Admin UI panel to edit tracked ships without editing JSON",
        "CSV export of position history per ship",
    ]
    for item in nxt:
        shape = slide.shapes.add_shape(9, Inches(0.4), y + Inches(0.04), Inches(0.14), Inches(0.14))
        shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
        add_textbox(slide, Inches(0.6), y, Inches(5.85), Inches(0.3),
                    item, font_size=10, color=GREY)
        y += Inches(0.31)

    # Right column – extensions
    add_textbox(slide, Inches(6.8), Inches(1.0), Inches(6.1), Inches(0.28),
                "Potential Extensions", font_size=12, bold=True, color=ACCENT)

    extensions = [
        "Alerting – notify when a ship enters a geofence or exceeds a speed threshold",
        "Weather overlay – add a second WebSocket API as a parallel Kafka stream; join with ship positions in PySpark",
        "Data warehouse – push from Kafka to Redshift / Snowflake for analytical queries",
        "AWS Glue + Lambda – serverless ETL for periodic aggregations",
        "Step Functions – orchestrate multi-step maintenance jobs",
        "User auth – JWT-gated access for the webapp",
        "Terraform – replace AWS CLI provisioning with declarative IaC",
    ]
    ey = Inches(1.32)
    for item in extensions:
        shape = slide.shapes.add_shape(9, Inches(6.8), ey + Inches(0.04), Inches(0.14), Inches(0.14))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x55)
        shape.line.fill.background()
        add_textbox(slide, Inches(7.0), ey, Inches(5.85), Inches(0.3),
                    item, font_size=10, color=GREY)
        ey += Inches(0.38)

    add_card(slide, Inches(6.8), Inches(4.6), Inches(6.1), Inches(2.1),
             title="Technologies demonstrated",
             bullets=[
                 "This project covers the full data engineering stack learned across 12 weeks:",
                 "Python · SQL · Docker · REST APIs · Kafka · PostgreSQL · Flask · MapLibre",
                 "AWS EC2 · GitHub Actions · CI/CD · Secrets Management · Stream Processing",
             ], font_size=11)


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)

    out = r"c:\Users\thoma\Desktop\DataEngineering\projects\myCapstoneProject\docs\AIS_Ship_Tracker_Capstone.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
